import streamlit as st
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_community.vectorstores import Chroma
from langchain_core.documents import Document
from langchain_core.prompts import ChatPromptTemplate
from langchain_ollama import ChatOllama, OllamaEmbeddings

import pandas as pd
import tempfile, os, subprocess, time

# ================= CONFIG =================
st.set_page_config(page_title="AI Đọc File Pro", layout="wide")

EMBED_MODEL = "nomic-embed-text"

# ================= GET MODELS =================
def get_models():
    try:
        result = subprocess.run(["ollama", "list"], capture_output=True, text=True)
        lines = result.stdout.split("\n")[1:]
        models = [l.split()[0] for l in lines if l.strip() and "embed" not in l]
        return models if models else ["qwen2.5:7b"]
    except:
        return ["qwen2.5:7b"]

# ================= SESSION =================
for key in ["messages", "documents", "vectorstore"]:
    if key not in st.session_state:
        st.session_state[key] = [] if key != "vectorstore" else None

# ================= SIDEBAR =================
with st.sidebar:
    st.title("⚙︎ Settings")

    models = get_models()
    selected_model = st.selectbox("✴ Model", models)

    temperature = st.slider("Temperature", 0.0, 1.0, 0.6)
    context_len = st.slider("Context Length", 2048, 16384, 8192)

    st.divider()

    st.write(f"🗁 Files: {len(st.session_state.documents)}")
    st.write(f"≫ Messages: {len(st.session_state.messages)}")

    if st.button("+ New Chat"):
        st.session_state.messages = []
        st.rerun()

    if st.button("🗑 Clear All"):
        st.session_state.messages = []
        st.session_state.documents = []
        st.session_state.vectorstore = None
        st.success("✔ Reset xong")
        st.rerun()

    st.caption("⬩➤ RAG Offline Pro")

# ================= MODEL =================
embeddings = OllamaEmbeddings(model=EMBED_MODEL)

llm = ChatOllama(
    model=selected_model,
    temperature=temperature,
    num_ctx=context_len
)

# ================= PARSER =================
def parse_file(path, name):
    name = name.lower()

    try:
        # PDF
        if name.endswith(".pdf"):
            from langchain_community.document_loaders import PyPDFLoader
            pages = PyPDFLoader(path).load()
            return "\n".join([p.page_content for p in pages])

        # WORD
        elif name.endswith(".docx"):
            import docx
            doc = docx.Document(path)
            text = [p.text for p in doc.paragraphs if p.text.strip()]

            for table in doc.tables:
                for row in table.rows:
                    row_text = " | ".join(cell.text for cell in row.cells if cell.text)
                    if row_text:
                        text.append(row_text)

            return "\n".join(text)

        # EXCEL
        elif name.endswith((".xlsx", ".xls", ".csv")):
            try:
                df = pd.read_excel(path, dtype=str)
            except:
                df = pd.read_csv(path, dtype=str)

            # giữ nguyên cấu trúc bảng
            df = df.fillna("")

            rows = []

            # header (giữ đủ cột)
            header = [str(col).strip() for col in df.columns]
            rows.append(" | ".join(header))

            # rows (KHÔNG được bỏ ô trống)
            for _, r in df.iterrows():
                row_values = []

                for cell in r:
                    try:
                        cell_str = str(cell).strip()
                    except:
                        cell_str = ""

                    # giữ placeholder để không lệch cột
                    if cell_str == "" or cell_str.lower() == "nan":
                        cell_str = "[EMPTY]"

                    row_values.append(cell_str)

                rows.append(" | ".join(row_values))

            return "\n".join(rows)

        # PPT
        elif name.endswith(".pptx"):
            from pptx import Presentation
            prs = Presentation(path)
            text = []

            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text.append(shape.text.strip())

            return "\n".join(text)

        # TXT
        else:
            return open(path, encoding="utf-8", errors="ignore").read()

    except Exception as e:
        return f"[ERROR]: {e}"

# ================= SPLIT (FIX CHÍNH) =================
def split_docs(docs):
    splitter = RecursiveCharacterTextSplitter(
        chunk_size=1000,        # lớn hơn
        chunk_overlap=300,      # overlap lớn → không mất đoạn cuối
        separators=["\n\n", "\n", ".", " "]
    )
    return splitter.split_documents(docs)

# ================= RETRIEVE (FIX CHÍNH) =================
def retrieve(query, vs):
    # lấy nhiều hơn
    docs = vs.similarity_search(query, k=12)
    return docs

# ================= CONTEXT (FIX CHÍNH) =================
def build_context(docs):
    context = ""
    for d in docs:
        context += f"\n\n📄 {d.metadata.get('source')}\n{d.page_content}"

    return context[:20000]  # tăng context

# ================= PROMPT =================
prompt = ChatPromptTemplate.from_template("""
Bạn là AI đọc tài liệu.

YÊU CẦU:
- Trả lời đầy đủ, không bỏ sót thông tin
- Nếu danh sách dài → liệt kê đầy đủ
- Không được tự ý cắt nội dung

Context:
{context}

Câu hỏi:
{question}
""")

# ================= UI =================
st.title("✴ AI Đọc File")

files = st.file_uploader("🗁 Upload file", accept_multiple_files=True)

# reset documents khi upload mới
if files:
    st.session_state.documents = []

    for f in files:
        with tempfile.NamedTemporaryFile(delete=False) as tmp:
            tmp.write(f.getvalue())
            path = tmp.name

        text = parse_file(path, f.name)

        st.session_state.documents.append(
            Document(page_content=text, metadata={"source": f.name})
        )

        os.unlink(path)

    st.session_state.vectorstore = None
    st.success("✔ Files loaded!")

# ================= VECTOR =================
if st.session_state.documents and st.session_state.vectorstore is None:
    with st.spinner("Indexing..."):
        splits = split_docs(st.session_state.documents)

        st.session_state.vectorstore = Chroma.from_documents(
            splits,
            embeddings,
            collection_name="rag_" + str(time.time())  # tránh reuse
        )

    st.success("⬩➤ Ready!")

# ================= CHAT =================
for i, msg in enumerate(st.session_state.messages):
    with st.chat_message(msg["role"]):
        st.markdown(msg["content"])

        if msg["role"] == "assistant":
            if st.button("✖", key=f"del_{i}"):
                st.session_state.messages.pop(i)
                st.rerun()

# ================= INPUT =================
if user_input := st.chat_input("≫ Hỏi về tài liệu..."):
    st.session_state.messages.append({"role": "user", "content": user_input})

    with st.chat_message("user"):
        st.markdown(user_input)

    with st.chat_message("assistant"):
        placeholder = st.empty()

        # loading
        for msg in ["Đang đọc tài liệu...", "Đang tìm thông tin...", "Đang tổng hợp..."]:
            placeholder.markdown("✴ " + msg)
            time.sleep(0.3)

        if not st.session_state.vectorstore:
            answer = "⚠️ Upload file trước."
            placeholder.markdown(answer)

        else:
            docs = retrieve(user_input, st.session_state.vectorstore)
            context = build_context(docs)

            chain = prompt | llm

            full = ""

            try:
                for chunk in chain.stream({
                    "context": context,
                    "question": user_input
                }):
                    token = chunk.content if hasattr(chunk, "content") else str(chunk)
                    full += token
                    placeholder.markdown(full + "▌")

            except Exception as e:
                full = f"✖ Lỗi: {str(e)}"

            placeholder.markdown(full)
            answer = full

    st.session_state.messages.append({"role": "assistant", "content": answer})
